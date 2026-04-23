from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Colors & Styles ---
    # Minimalist color palette
    NAVY = RGBColor(10, 25, 47)
    BLUE = RGBColor(0, 112, 243)
    GRAY = RGBColor(102, 102, 102)
    WHITE = RGBColor(255, 255, 255)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_minimal_title(slide, title_text, subtitle_text=None):
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = NAVY
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        if subtitle_text:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = GRAY
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Placeholder logo (Rectangle)
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(1), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.color.rgb = BLUE

    add_minimal_title(slide, "MarketMind", "AI-Orchestrated Financial Research Platform")

    # --- Slide 2: Overview ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    add_minimal_title(slide, "Project Overview")
    
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "MarketMind is a multi-agent system designed for the Indian stock and mutual fund markets."
    p.level = 0
    
    bullet_points = [
        "Synthesizes quantitative metrics, news sentiment, and historical trends.",
        "Provides retail-investor-friendly insights via AI agents.",
        "Interactive canvas for head-to-head comparisons and risk analysis.",
        "Automated data pipeline using GitHub Actions."
    ]
    
    for point in bullet_points:
        p = content.add_paragraph()
        p.text = point
        p.level = 1

    # --- Slide 3: Tech Stack ---
    slide = prs.slides.add_slide(slide_layout)
    add_minimal_title(slide, "The Tech Stack")
    
    content = slide.placeholders[1].text_frame
    
    # Frontend
    p = content.add_paragraph()
    p.text = "Frontend: Next.js 15 (App Router)"
    p.font.bold = True
    p.level = 0
    for item in ["TypeScript & Tailwind CSS", "Recharts for visualization", "Zustand State Management"]:
        p = content.add_paragraph()
        p.text = item
        p.level = 1

    # Backend
    p = content.add_paragraph()
    p.text = "Backend: FastAPI (Python)"
    p.font.bold = True
    p.level = 0
    for item in ["Groq AI (Llama 3.1 70b)", "YFinance & MFAPI.in Data", "Supabase (PostgreSQL)"]:
        p = content.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: AI Agent Pipeline ---
    slide = prs.slides.add_slide(slide_layout)
    add_minimal_title(slide, "Multi-Agent Architecture")
    
    content = slide.placeholders[1].text_frame
    agents = [
        ("Router Agent", "Classifies user intent (Quant, News, Screener)."),
        ("Quant Agent", "Fetches real-time market data & computes metrics."),
        ("News Parser", "Scrapes RSS feeds and assigns AI-driven sentiment."),
        ("Synthesis Core", "Combines data into structured, actionable insights.")
    ]
    
    for title, desc in agents:
        p = content.add_paragraph()
        p.text = f"{title}: {desc}"
        p.level = 0

    # --- Slide 5: Interactive Canvas ---
    slide = prs.slides.add_slide(slide_layout)
    add_minimal_title(slide, "Interactive Canvas UI")
    
    content = slide.placeholders[1].text_frame
    features = [
        "Split-View Layout: Chat conversation on left, deep-dive Canvas on right.",
        "MF Comparison: Normalized NAV charts for head-to-head analysis.",
        "Risk Metrics: Automated Alpha, Beta, Sharpe, and CAGR calculation.",
        "Portfolio Overlap: Concentration analysis and shared holdings."
    ]
    
    for feature in features:
        p = content.add_paragraph()
        p.text = feature
        p.level = 1

    # --- Slide 6: Data & Automation ---
    slide = prs.slides.add_slide(slide_layout)
    add_minimal_title(slide, "Data Engine & Automation")
    
    content = slide.placeholders[1].text_frame
    p = content.add_paragraph()
    p.text = "Scheduled Pipelines (GitHub Actions):"
    p.font.bold = True
    p.level = 0
    
    p = content.add_paragraph()
    p.text = "EOD Stock Fetch: Daily 16:30 IST update for Nifty 50/100/250."
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "MF Sync: Automated metadata synchronization with Supabase."
    p.level = 1

    p = content.add_paragraph()
    p.text = "Data Sources:"
    p.font.bold = True
    p.level = 0
    p = content.add_paragraph()
    p.text = "YFinance, Google News RSS, and MFAPI.in."
    p.level = 1

    # --- Slide 7: Deployment Architecture ---
    slide = prs.slides.add_slide(slide_layout)
    add_minimal_title(slide, "Deployment & Infrastructure")
    
    content = slide.placeholders[1].text_frame
    infra = [
        ("Frontend", "Vercel (Edge Runtime ready)"),
        ("Backend", "Render (Python FastAPI)"),
        ("Database", "Supabase (Cloud PostgreSQL)"),
        ("Automation", "GitHub Actions (CI/CD + Data Cron)")
    ]
    
    for part, service in infra:
        p = content.add_paragraph()
        p.text = f"{part} deployed on {service}"
        p.level = 0

    # --- Slide 8: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_minimal_title(slide, "Thank You", "MarketMind: Empowering Retail Investors")
    
    # Save the presentation
    prs.save('MarketMind_Project_Presentation.pptx')
    print("Presentation generated successfully: MarketMind_Project_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
